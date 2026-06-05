"""
Genera plantilla PowerPoint para manual de usuario NOM-035
Colores: sistema Intra NOM-035 (variables.css)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paleta ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x15, 0x26)   # Azul noche profundo
BG2      = RGBColor(0x13, 0x1D, 0x33)   # Superficie oscura
ACCENT   = RGBColor(0x03, 0xC4, 0xCE)   # Cyan marca
ACCENT_D = RGBColor(0x02, 0x7F, 0x88)   # Cyan oscuro
TEXT1    = RGBColor(0xE4, 0xED, 0xF6)   # Texto primario
TEXT2    = RGBColor(0xA4, 0xBB, 0xCF)   # Texto secundario
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SCRBG    = RGBColor(0xEE, 0xF2, 0xFA)   # Fondo de capturas (mismo que app)
BORDER   = RGBColor(0x1E, 0x30, 0x4F)   # Borde oscuro
G_OK     = RGBColor(0x10, 0xB9, 0x81)   # Verde éxito
G_WARN   = RGBColor(0xF5, 0x9E, 0x0B)   # Ámbar advertencia
G_INFO   = RGBColor(0x03, 0xC4, 0xCE)   # Info = accent
G_ERR    = RGBColor(0xEF, 0x44, 0x44)   # Rojo importante

W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────────────
def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def add_rect(slide, left, top, width, height, fill=None, border=None, bw=1):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(bw)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=TEXT1,
                align=PP_ALIGN.LEFT, wrap=True,
                italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bar(slide, left=Inches(0), top=Inches(0),
                   width=Inches(0.06), height=H, color=ACCENT):
    """Barra vertical de acento."""
    add_rect(slide, left, top, width, height, fill=color)

def add_step_circle(slide, number, left, top, size=Inches(0.45), color=ACCENT):
    """Círculo numerado para pasos."""
    shape = slide.shapes.add_shape(9, left, top, size, size)  # oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = str(number)
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    return shape

def add_screenshot_placeholder(slide, left, top, width, height,
                                label="Insertar captura de pantalla aquí"):
    """Rectángulo blanco/grisáceo para captura — contrasta contra fondo oscuro."""
    # Fondo igual al app (EEF2FA) para que las capturas se fundan
    body = add_rect(slide, left, top, width, height,
                    fill=SCRBG, border=ACCENT, bw=1.5)

    # Texto guía centrado
    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), top + height // 2 - Pt(20),
        width - Inches(0.2), Pt(40)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x6E, 0x8B, 0xAA)
    run.font.italic = True
    run.font.name = "Calibri"
    return body

def add_tag(slide, text, left, top, color=ACCENT):
    """Etiqueta pequeña de módulo."""
    w = Inches(1.6)
    h = Inches(0.28)
    add_rect(slide, left, top, w, h,
             fill=RGBColor(int(color[0]*0.15), int(color[1]*0.15), int(color[2]*0.15)),
             border=color, bw=0.75)
    add_textbox(slide, text,
                left + Inches(0.08), top + Inches(0.02),
                w - Inches(0.16), h,
                font_size=8, bold=True, color=color, align=PP_ALIGN.LEFT)

def add_page_number(slide, num, total):
    add_textbox(slide, f"{num}  /  {total}",
                W - Inches(1.2), H - Inches(0.4),
                Inches(1.0), Inches(0.3),
                font_size=9, color=TEXT2, align=PP_ALIGN.RIGHT)

def add_divider_line(slide, top, color=BORDER):
    add_rect(slide, Inches(0.5), top, W - Inches(1.0), Emu(18000),
             fill=color)

# ── Presentación ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # completamente en blanco

TOTAL_SLIDES = 9

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)

# Barra izquierda cyan
add_rect(s, 0, 0, Inches(0.45), H, fill=ACCENT)

# Barra cyan inferior delgada
add_rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)

# Bloque de marca — fondo ligeramente más claro
add_rect(s, Inches(0.45), 0, W - Inches(0.45), H, fill=BG)

# Caja decorativa superior derecha
add_rect(s, W - Inches(3.5), Inches(0), Inches(3.5), Inches(1.8),
         fill=BG2)

# Norma label
add_textbox(s, "NOM-035-STPS-2018",
            W - Inches(3.4), Inches(0.35),
            Inches(3.3), Inches(0.5),
            font_size=11, bold=True, color=ACCENT,
            align=PP_ALIGN.RIGHT)

add_textbox(s, "Factores de Riesgo Psicosocial",
            W - Inches(3.4), Inches(0.75),
            Inches(3.3), Inches(0.4),
            font_size=9, color=TEXT2, align=PP_ALIGN.RIGHT)

# Línea separadora vertical decorativa
add_rect(s, Inches(0.75), Inches(2.2), Emu(18000), Inches(3.1),
         fill=ACCENT)

# Título principal
add_textbox(s, "Manual de Usuario",
            Inches(1.0), Inches(2.0),
            Inches(8.5), Inches(1.0),
            font_size=40, bold=True, color=TEXT1,
            align=PP_ALIGN.LEFT, font_name="Calibri Light")

add_textbox(s, "Plataforma Intra NOM-035",
            Inches(1.0), Inches(2.95),
            Inches(8.0), Inches(0.65),
            font_size=26, bold=False, color=ACCENT,
            align=PP_ALIGN.LEFT, font_name="Calibri Light")

add_textbox(s, "Guía de operación para administradores de empresa (tenants)",
            Inches(1.0), Inches(3.65),
            Inches(8.5), Inches(0.55),
            font_size=14, color=TEXT2, align=PP_ALIGN.LEFT)

# Separador
add_divider_line(s, Inches(4.5))

# Bloque inferior: versión / fecha
add_textbox(s, "Versión 1.0",
            Inches(1.0), Inches(4.7),
            Inches(2.5), Inches(0.35),
            font_size=11, bold=True, color=TEXT2, align=PP_ALIGN.LEFT)
add_textbox(s, "© 2025 · Uso interno",
            W - Inches(3.5), Inches(4.7),
            Inches(3.3), Inches(0.35),
            font_size=11, color=TEXT2, align=PP_ALIGN.RIGHT)

# Caja placeholder logo
add_rect(s, Inches(1.0), Inches(5.3), Inches(2.8), Inches(1.2),
         fill=BG2, border=BORDER, bw=1)
add_textbox(s, "[ Logo de la empresa ]",
            Inches(1.0), Inches(5.7),
            Inches(2.8), Inches(0.4),
            font_size=9, color=TEXT2,
            align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ÍNDICE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)
add_accent_bar(s)
add_rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)
add_page_number(s, 2, TOTAL_SLIDES)

add_textbox(s, "Contenido del manual",
            Inches(0.55), Inches(0.35),
            Inches(9.0), Inches(0.6),
            font_size=26, bold=True, color=TEXT1,
            align=PP_ALIGN.LEFT, font_name="Calibri Light")
add_divider_line(s, Inches(1.05))

modulos = [
    ("01", "Acceso y configuración inicial",    "Inicio de sesión, perfil de empresa y ciclos NOM"),
    ("02", "Gestión de trabajadores",            "Alta, importación desde Excel y datos demográficos"),
    ("03", "Cuestionarios",                      "Asignación, envío de ligas y seguimiento de respuestas"),
    ("04", "Resultados y diagnóstico",           "Dashboard analítico, niveles de riesgo y gráficas"),
    ("05", "Plan de acción",                     "Registro de medidas correctivas por módulo"),
    ("06", "Política y documentos",              "Carga de política de prevención y evidencias"),
    ("07", "Informe NOM-035",                    "Generación de PDF oficial para IMSS/STPS"),
]

col1 = modulos[:4]
col2 = modulos[4:]

for i, (num, titulo, desc) in enumerate(col1):
    top = Inches(1.3) + i * Inches(1.35)
    # Número
    add_step_circle(s, num, Inches(0.6), top + Inches(0.05),
                    size=Inches(0.42), color=ACCENT)
    add_textbox(s, titulo,
                Inches(1.15), top,
                Inches(5.4), Inches(0.4),
                font_size=13, bold=True, color=TEXT1)
    add_textbox(s, desc,
                Inches(1.15), top + Inches(0.38),
                Inches(5.4), Inches(0.35),
                font_size=10, color=TEXT2)

for i, (num, titulo, desc) in enumerate(col2):
    top = Inches(1.3) + i * Inches(1.35)
    add_step_circle(s, num, Inches(7.15), top + Inches(0.05),
                    size=Inches(0.42), color=ACCENT)
    add_textbox(s, titulo,
                Inches(7.7), top,
                Inches(5.4), Inches(0.4),
                font_size=13, bold=True, color=TEXT1)
    add_textbox(s, desc,
                Inches(7.7), top + Inches(0.38),
                Inches(5.4), Inches(0.35),
                font_size=10, color=TEXT2)

# Línea divisora entre columnas
add_rect(s, Inches(6.9), Inches(1.15), Emu(18000), Inches(5.7), fill=BORDER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SEPARADOR DE MÓDULO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)

# Bloque izquierdo de color sólido (30% del ancho)
add_rect(s, 0, 0, Inches(3.9), H, fill=ACCENT)

# Número grande del módulo (muy tenue, decorativo)
add_textbox(s, "01",
            Inches(0.1), Inches(0.5),
            Inches(3.7), Inches(5.0),
            font_size=180, bold=True,
            color=RGBColor(0x02, 0x9A, 0xA2),
            align=PP_ALIGN.CENTER, font_name="Calibri Light")

# Título del módulo
add_textbox(s, "MÓDULO 01",
            Inches(4.3), Inches(2.2),
            Inches(8.5), Inches(0.5),
            font_size=13, bold=True, color=ACCENT,
            align=PP_ALIGN.LEFT)

add_textbox(s, "Acceso y configuración inicial",
            Inches(4.3), Inches(2.75),
            Inches(8.7), Inches(0.85),
            font_size=34, bold=True, color=TEXT1,
            align=PP_ALIGN.LEFT, font_name="Calibri Light")

add_divider_line(s, Inches(3.75), color=BORDER)

add_textbox(s, "En esta sección aprenderás a:",
            Inches(4.3), Inches(4.0),
            Inches(8.5), Inches(0.4),
            font_size=12, color=TEXT2, align=PP_ALIGN.LEFT)

for i, item in enumerate([
    "Iniciar sesión en la plataforma",
    "Configurar el perfil de tu empresa",
    "Crear y gestionar ciclos NOM-035",
]):
    top = Inches(4.5) + i * Inches(0.45)
    add_rect(s, Inches(4.3), top + Inches(0.13),
             Inches(0.07), Inches(0.18), fill=ACCENT)
    add_textbox(s, item,
                Inches(4.5), top,
                Inches(8.5), Inches(0.42),
                font_size=12, color=TEXT1)

add_page_number(s, 3, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TEXTO IZQUIERDA + CAPTURA DERECHA (layout principal)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)
add_accent_bar(s)
add_rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)
add_page_number(s, 4, TOTAL_SLIDES)

add_tag(s, "MÓDULO 01", Inches(0.55), Inches(0.32))

add_textbox(s, "Inicio de sesión",
            Inches(0.55), Inches(0.72),
            Inches(5.0), Inches(0.55),
            font_size=24, bold=True, color=TEXT1,
            font_name="Calibri Light")

add_divider_line(s, Inches(1.35))

pasos = [
    "Ingresa la URL de la plataforma en tu navegador.",
    "Escribe tu usuario y contraseña asignados.",
    "Haz clic en Iniciar sesión.",
    "Si olvidaste tu contraseña, usa el enlace\n'Recuperar acceso' en la pantalla de login.",
]
for i, paso in enumerate(pasos):
    top = Inches(1.55) + i * Inches(1.15)
    add_step_circle(s, i + 1, Inches(0.55), top, size=Inches(0.38))
    add_textbox(s, paso,
                Inches(1.08), top - Inches(0.04),
                Inches(5.0), Inches(0.85),
                font_size=12, color=TEXT1)

# Nota al pie
add_rect(s, Inches(0.55), Inches(6.3), Inches(5.5), Inches(0.75),
         fill=RGBColor(0x13, 0x1D, 0x33), border=G_INFO, bw=1)
add_textbox(s, "ⓘ  Tu administrador de sistema te proporcionará las credenciales iniciales.",
            Inches(0.7), Inches(6.4),
            Inches(5.2), Inches(0.55),
            font_size=10, color=TEXT2, italic=True)

# Captura derecha
add_screenshot_placeholder(s,
    left=Inches(6.2), top=Inches(0.35),
    width=Inches(6.85), height=Inches(6.75))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CAPTURA PANTALLA COMPLETA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)
add_accent_bar(s)
add_rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)
add_page_number(s, 5, TOTAL_SLIDES)

add_tag(s, "MÓDULO 04", Inches(0.55), Inches(0.28))

add_textbox(s, "Dashboard de Resultados",
            Inches(0.55), Inches(0.62),
            Inches(10.0), Inches(0.55),
            font_size=22, bold=True, color=TEXT1,
            font_name="Calibri Light")

add_textbox(s, "Vista general del diagnóstico analítico NOM-035 por ciclo.",
            Inches(0.55), Inches(1.18),
            Inches(12.5), Inches(0.38),
            font_size=12, color=TEXT2)

add_divider_line(s, Inches(1.6))

# Captura grande
add_screenshot_placeholder(s,
    left=Inches(0.55), top=Inches(1.78),
    width=Inches(12.23), height=Inches(5.3),
    label="Insertar captura del Dashboard de Resultados aquí")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PASOS NUMERADOS (flujo de proceso)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)
add_accent_bar(s)
add_rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)
add_page_number(s, 6, TOTAL_SLIDES)

add_tag(s, "MÓDULO 03", Inches(0.55), Inches(0.28))

add_textbox(s, "Flujo completo: cuestionarios",
            Inches(0.55), Inches(0.62),
            Inches(10.0), Inches(0.55),
            font_size=22, bold=True, color=TEXT1,
            font_name="Calibri Light")
add_divider_line(s, Inches(1.25))

flujo = [
    (ACCENT,  "Crear ciclo NOM",
     "Ve a Onboarding → Ciclos. Crea un nuevo ciclo con el año\ny los datos de la empresa."),
    (G_OK,    "Asignar cuestionarios",
     "En Cuestionarios → Aplicaciones, asigna Guía I, III\ny/o V a cada trabajador."),
    (G_WARN,  "Enviar ligas de acceso",
     "El sistema genera una URL única por trabajador.\nCópiala y envíala por correo o WhatsApp."),
    (G_ERR,   "Monitorear avance",
     "El semáforo de estado muestra: Pendiente / En proceso\n/ Completado en tiempo real."),
]

col_w = Inches(3.0)
for i, (color, titulo, desc) in enumerate(flujo):
    left = Inches(0.55) + i * (col_w + Inches(0.18))
    # Barra de color superior
    add_rect(s, left, Inches(1.45), col_w, Inches(0.06), fill=color)
    # Tarjeta
    add_rect(s, left, Inches(1.51), col_w, Inches(4.6),
             fill=BG2, border=BORDER, bw=1)
    # Número
    add_step_circle(s, i + 1, left + Inches(0.15), Inches(1.7),
                    size=Inches(0.5), color=color)
    # Título
    add_textbox(s, titulo,
                left + Inches(0.1), Inches(2.35),
                col_w - Inches(0.2), Inches(0.45),
                font_size=13, bold=True, color=TEXT1)
    # Descripción
    add_textbox(s, desc,
                left + Inches(0.1), Inches(2.85),
                col_w - Inches(0.2), Inches(1.1),
                font_size=10, color=TEXT2)
    # Captura mini
    add_screenshot_placeholder(s,
        left=left + Inches(0.1),
        top=Inches(4.1),
        width=col_w - Inches(0.2),
        height=Inches(1.7),
        label="Captura")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DOS CAPTURAS LADO A LADO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)
add_accent_bar(s)
add_rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)
add_page_number(s, 7, TOTAL_SLIDES)

add_tag(s, "MÓDULO 02", Inches(0.55), Inches(0.28))

add_textbox(s, "Importación de trabajadores",
            Inches(0.55), Inches(0.62),
            Inches(10.0), Inches(0.55),
            font_size=22, bold=True, color=TEXT1,
            font_name="Calibri Light")
add_divider_line(s, Inches(1.25))

# Sub-título izquierdo
add_textbox(s, "① Descargar plantilla Excel",
            Inches(0.55), Inches(1.4),
            Inches(6.0), Inches(0.4),
            font_size=13, bold=True, color=ACCENT)
add_textbox(s, "Descarga la plantilla desde el botón\n\"Descargar plantilla\" y completa los datos.",
            Inches(0.55), Inches(1.82),
            Inches(6.0), Inches(0.65),
            font_size=11, color=TEXT2)

add_screenshot_placeholder(s,
    left=Inches(0.55), top=Inches(2.5),
    width=Inches(5.9), height=Inches(4.6),
    label="Captura: botón exportar / plantilla Excel")

# Sub-título derecho
add_textbox(s, "② Importar el archivo completado",
            Inches(6.85), Inches(1.4),
            Inches(6.0), Inches(0.4),
            font_size=13, bold=True, color=ACCENT)
add_textbox(s, "Usa el botón \"Importar Excel\" para\ncargar masivamente los trabajadores.",
            Inches(6.85), Inches(1.82),
            Inches(6.0), Inches(0.65),
            font_size=11, color=TEXT2)

add_screenshot_placeholder(s,
    left=Inches(6.85), top=Inches(2.5),
    width=Inches(5.9), height=Inches(4.6),
    label="Captura: modal de importación / resultado")

# Divisor vertical
add_rect(s, Inches(6.6), Inches(1.35), Emu(18000), Inches(5.85), fill=BORDER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CALLOUT / NOTA IMPORTANTE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)
add_accent_bar(s)
add_rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)
add_page_number(s, 8, TOTAL_SLIDES)

add_tag(s, "MÓDULO 04", Inches(0.55), Inches(0.28))
add_textbox(s, "Cálculo del diagnóstico",
            Inches(0.55), Inches(0.62),
            Inches(9.0), Inches(0.55),
            font_size=22, bold=True, color=TEXT1,
            font_name="Calibri Light")
add_divider_line(s, Inches(1.25))

# Captura principal
add_screenshot_placeholder(s,
    left=Inches(0.55), top=Inches(1.45),
    width=Inches(7.5), height=Inches(4.4))

# Callouts apilados a la derecha
callouts = [
    (G_INFO,  "ℹ Antes de calcular",
     "Asegúrate de que todos los cuestionarios del\nciclo estén en estado \"Completado\"."),
    (G_WARN,  "⚠ Recalcular",
     "Puedes recalcular cuantas veces necesites.\nLos resultados anteriores se sobrescriben."),
    (G_OK,    "✓ Resultado inmediato",
     "El dashboard se actualiza al instante\ndespués de presionar 'Calcular diagnóstico'."),
]
for i, (color, titulo, texto) in enumerate(callouts):
    top = Inches(1.45) + i * Inches(1.58)
    add_rect(s, Inches(8.35), top,
             Inches(4.7), Inches(1.42),
             fill=BG2, border=color, bw=1.5)
    # Barra lateral izquierda de color
    add_rect(s, Inches(8.35), top,
             Inches(0.06), Inches(1.42), fill=color)
    add_textbox(s, titulo,
                Inches(8.55), top + Inches(0.18),
                Inches(4.4), Inches(0.38),
                font_size=12, bold=True, color=color)
    add_textbox(s, texto,
                Inches(8.55), top + Inches(0.58),
                Inches(4.4), Inches(0.75),
                font_size=10, color=TEXT2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CIERRE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG)

# Faja izquierda cyan
add_rect(s, 0, 0, Inches(0.45), H, fill=ACCENT)
add_rect(s, 0, H - Inches(0.06), W, Inches(0.06), fill=ACCENT)

# Bloque decorativo inferior
add_rect(s, Inches(0.45), H - Inches(2.2), W - Inches(0.45), Inches(2.14),
         fill=BG2)

add_textbox(s, "¿Dudas o soporte?",
            Inches(1.0), Inches(1.8),
            Inches(11.0), Inches(0.7),
            font_size=36, bold=True, color=TEXT1,
            align=PP_ALIGN.CENTER, font_name="Calibri Light")

add_textbox(s, "Contacta a tu administrador de sistema o al equipo de soporte técnico.",
            Inches(1.5), Inches(2.6),
            Inches(10.0), Inches(0.5),
            font_size=14, color=TEXT2, align=PP_ALIGN.CENTER)

# Separador
add_divider_line(s, Inches(3.3))

# Datos de contacto (placeholders)
contactos = [
    ("✉  Correo",   "soporte@tuempresa.com"),
    ("☎  Teléfono", "+52 (55) 0000-0000"),
    ("🌐  Plataforma", "https://tudominio.com"),
]
for i, (label, valor) in enumerate(contactos):
    left = Inches(1.3) + i * Inches(3.9)
    add_rect(s, left, Inches(3.6), Inches(3.5), Inches(1.3),
             fill=BG2, border=BORDER, bw=1)
    add_textbox(s, label,
                left + Inches(0.15), Inches(3.72),
                Inches(3.2), Inches(0.38),
                font_size=10, bold=True, color=ACCENT)
    add_textbox(s, valor,
                left + Inches(0.15), Inches(4.1),
                Inches(3.2), Inches(0.4),
                font_size=12, color=TEXT1)

add_textbox(s, "NOM-035-STPS-2018  ·  Manual de Usuario  ·  v1.0",
            Inches(0.7), H - Inches(1.7),
            W - Inches(1.4), Inches(0.4),
            font_size=10, color=TEXT2,
            align=PP_ALIGN.CENTER)

add_textbox(s, "© 2025 — Uso exclusivo para personal autorizado",
            Inches(0.7), H - Inches(1.3),
            W - Inches(1.4), Inches(0.4),
            font_size=9, color=RGBColor(0x53, 0x6F, 0x8A),
            align=PP_ALIGN.CENTER)


# ── Guardar ──────────────────────────────────────────────────────────────────
out = "Manual_NOM035_Plantilla.pptx"
prs.save(out)
print(f"Archivo guardado: {out}")
